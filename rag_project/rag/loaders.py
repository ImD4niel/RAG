"""
loaders.py — Multi-format document ingestion for DRAGON RAG.

Key fixes vs. previous version:
- load_audio_video() now correctly awaits the async OpenAI Whisper API call.
  Previously the coroutine was called without `await`, causing a runtime error
  on any audio/video file.
- All print() calls replaced with structured logger calls so that errors appear
  in the unified rag.log rather than being silently lost in stdout.
- A configurable file-size guard prevents runaway memory or compute usage when
  very large files are placed in the data directory.
- The media_dir path is computed once via a helper function instead of being
  repeated (fragile) in every loader.
- Full type hints and docstrings.
"""

from __future__ import annotations

import logging
import os
from io import BytesIO

import requests
from bs4 import BeautifulSoup

from .cleaner import clean_text
from .config import get_settings

logger = logging.getLogger(__name__)




_DEFAULT_MAX_FILE_BYTES = 100 * 1024 * 1024  





def _get_media_dir() -> str:
    """
    Return the absolute path to the shared media extraction directory.

    All loaders write extracted images to this directory so that the pipeline
    can serve them during multimodal generation.
    """
    base = os.path.dirname(  
        os.path.dirname(       
            os.path.abspath(__file__)
        )
    )
    media_dir = os.path.join(base, "..", "data", "media")
    os.makedirs(media_dir, exist_ok=True)
    return os.path.normpath(media_dir)


def _file_too_large(filepath: str, max_bytes: int = _DEFAULT_MAX_FILE_BYTES) -> bool:
    """Return True if the file exceeds the configured size limit."""
    try:
        size = os.path.getsize(filepath)
        if size > max_bytes:
            logger.warning(
                f"Skipping {os.path.basename(filepath)}: "
                f"file size {size / 1024 / 1024:.1f} MB exceeds "
                f"limit of {max_bytes / 1024 / 1024:.0f} MB."
            )
            return True
        return False
    except OSError as e:
        logger.error(f"Cannot stat {filepath}: {e}")
        return True





def extract_text_with_ocr(pdf_path: str) -> list[dict]:
    """Extract text from a PDF using Tesseract OCR (page-by-page)."""
    try:
        import pytesseract
        from pdf2image import convert_from_path
    except ImportError as e:
        logger.warning(f"OCR dependencies missing ({e}) — skipping {pdf_path}")
        return []

    images = convert_from_path(pdf_path)
    docs = []
    for i, image in enumerate(images):
        text = pytesseract.image_to_string(image)
        cleaned = clean_text(text)
        if cleaned:
            docs.append(
                {"text": cleaned, "source": pdf_path, "page": i + 1, "doc_type": "pdf_ocr"}
            )
    return docs





async def extract_text_with_vlm(image_pil: "Image.Image") -> str:  # type: ignore[name-defined]
    """
    Describe an image using the configured vision LLM.

    NOTE: The active ingestion pipeline does NOT call this function directly.
    Image chunks are flagged with `needs_vlm_description=True` during load_pdf()
    and the async description + embedding happens in vector_store._embed_chunks_with_semaphore().
    This function is a convenience wrapper for direct callers only.
    """
    from .llm import describe_image
    buffer = BytesIO()
    image_pil.save(buffer, format="JPEG")
    description = await describe_image(buffer.getvalue())
    return f"[IMAGE DESCRIPTION: {description}]"





def load_txt(filepath: str) -> list[dict]:
    """Load a plain-text file and return as a single chunk."""
    if _file_too_large(filepath):
        return []
    try:
        with open(filepath, "r", encoding="utf-8", errors="replace") as f:
            cleaned = clean_text(f.read())
        if not cleaned:
            return []
        return [{"text": cleaned, "source": filepath, "page": None, "doc_type": "text"}]
    except OSError as e:
        logger.error(f"Failed to load TXT {filepath}: {e}")
        return []


def _tables_to_markdown(pdfplumber_page: object) -> list[str]:
    """
    Extract tables from a pdfplumber page and convert each to Markdown.

    Returns a list of Markdown table strings. Each table is atomic — it is
    never split across chunk boundaries during chunking.
    """
    markdown_tables = []
    try:
        tables = pdfplumber_page.extract_tables()
        if not tables:
            return []
        for table in tables:
            if not table:
                continue
            rows = []
            for r_idx, row in enumerate(table):
                cells = [str(cell).strip() if cell is not None else "" for cell in row]
                rows.append("| " + " | ".join(cells) + " |")
                if r_idx == 0:
                    rows.append("|" + "|".join([" --- "] * len(cells)) + "|")
            markdown_tables.append("\n".join(rows))
    except Exception as e:
        logger.warning(f"pdfplumber table extraction failed: {e}")
    return markdown_tables


def load_pdf(filepath: str) -> list[dict]:
    """
    Load a PDF using PyMuPDF (text) + pdfplumber (tables) + media extraction (images).

    Pipeline:
    1. pdfplumber extracts tables as Markdown (preserving row/column structure).
    2. PyMuPDF extracts prose text page-by-page (fast, handles Unicode well).
    3. Images are extracted via PyMuPDF, saved to media/, and flagged with
       `needs_vlm_description=True` so the embedding pipeline can call the
       vision LLM asynchronously during ingestion.

    Tables are stored as a single atomic chunk (doc_type='pdf_table') and are
    never split during chunking — their structure must remain intact.
    """
    if _file_too_large(filepath):
        return []
    try:
        import fitz
    except ImportError:
        logger.error("PyMuPDF (fitz) not installed — cannot load PDF.")
        return []

    try:
        doc = fitz.open(filepath)
    except Exception as e:
        logger.error(f"Failed to open PDF {filepath}: {e}")
        return []

    docs: list[dict] = []
    media_dir = _get_media_dir()
    base_name = os.path.splitext(os.path.basename(filepath))[0]

    page_table_texts: dict[int, set[str]] = {}  
    try:
        import pdfplumber
        with pdfplumber.open(filepath) as pl_pdf:
            for page_num, pl_page in enumerate(pl_pdf.pages):
                md_tables = _tables_to_markdown(pl_page)
                page_table_texts[page_num] = set()
                for md_table in md_tables:
                    if md_table.strip():
                        docs.append({
                            "text": md_table,
                            "source": filepath,
                            "page": page_num + 1,
                            "doc_type": "pdf_table",
                            "media_type": "text",
                            "media_path": None,
                            "structure_type": "table",
                        })
                        for cell_text in md_table.replace("|", " ").split():
                            page_table_texts[page_num].add(cell_text.lower())
    except ImportError:
        logger.warning("pdfplumber not installed — tables will be extracted as flat prose. "
                       "Install with: pip install pdfplumber")
    except Exception as e:
        logger.warning(f"pdfplumber failed for {filepath}: {e} — falling back to prose-only.")

    for i, page in enumerate(doc):
        text = page.get_text()
        if text:
            cleaned = clean_text(text)
            if cleaned:
                docs.append({
                    "text": cleaned,
                    "source": filepath,
                    "page": i + 1,
                    "doc_type": "pdf",
                    "media_type": "text",
                    "media_path": None,
                    "structure_type": "paragraph",
                })

        for img_index, img in enumerate(page.get_images(full=True)):
            xref = img[0]
            try:
                base_image = doc.extract_image(xref)
                img_ext = base_image.get("ext", "png")
                img_filename = f"{base_name}_page{i+1}_img{img_index}.{img_ext}"
                img_path = os.path.join(media_dir, img_filename)
                with open(img_path, "wb") as img_file:
                    img_file.write(base_image["image"])

                docs.append({
                    "text": f"[Image on page {i+1} from {base_name} — pending VLM description]",
                    "source": filepath,
                    "page": i + 1,
                    "doc_type": "pdf_image",
                    "media_type": "image",
                    "media_path": img_path,
                    "needs_vlm_description": True,
                })
            except Exception as e:
                logger.warning(f"Failed to extract image xref={xref} from {filepath}: {e}")

    return docs


def load_docx(filepath: str) -> list[dict]:
    """Load a DOCX file — extracts body text and embedded images."""
    if _file_too_large(filepath):
        return []
    try:
        from docx import Document
    except ImportError:
        logger.error("python-docx not installed — cannot load DOCX.")
        return []

    try:
        doc = Document(filepath)
    except Exception as e:
        logger.error(f"Failed to open DOCX {filepath}: {e}")
        return []

    docs: list[dict] = []
    media_dir = _get_media_dir()
    base_name = os.path.splitext(os.path.basename(filepath))[0]

    text = "\n".join(p.text for p in doc.paragraphs)
    cleaned = clean_text(text)
    if cleaned:
        docs.append(
            {
                "text": cleaned,
                "source": filepath,
                "page": None,
                "doc_type": "docx",
                "media_type": "text",
                "media_path": None,
            }
        )

    for rel in doc.part.rels.values():
        if "image" not in rel.target_ref:
            continue
        try:
            img_blob = rel.target_part.blob
            img_ext = rel.target_part.content_type.split("/")[-1]
            if not img_ext or len(img_ext) > 5:
                img_ext = "png"
            img_path = os.path.join(media_dir, f"{base_name}_img_{rel.rId}.{img_ext}")
            with open(img_path, "wb") as f:
                f.write(img_blob)
            docs.append(
                {
                    "text": f"[Image {rel.rId} from {base_name}]",
                    "source": filepath,
                    "page": None,
                    "doc_type": "docx_image",
                    "media_type": "image",
                    "media_path": img_path,
                }
            )
        except Exception as e:
            logger.warning(f"Failed to extract DOCX image {rel.rId}: {e}")

    return docs


def load_pptx(filepath: str) -> list[dict]:
    """Load a PPTX file — extracts slide text and embedded images."""
    if _file_too_large(filepath):
        return []
    try:
        from pptx import Presentation
    except ImportError:
        logger.warning("python-pptx not installed — skipping PPTX.")
        return []

    try:
        prs = Presentation(filepath)
    except Exception as e:
        logger.error(f"Failed to open PPTX {filepath}: {e}")
        return []

    docs: list[dict] = []
    media_dir = _get_media_dir()
    base_name = os.path.splitext(os.path.basename(filepath))[0]

    for i, slide in enumerate(prs.slides):
        slide_text_parts: list[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                slide_text_parts.append(shape.text)
            if hasattr(shape, "image"):
                try:
                    img_path = os.path.join(
                        media_dir,
                        f"{base_name}_slide{i+1}_img_{shape.shape_id}.{shape.image.ext}",
                    )
                    with open(img_path, "wb") as f:
                        f.write(shape.image.blob)
                    docs.append(
                        {
                            "text": f"[Image on slide {i+1} from {base_name}]",
                            "source": filepath,
                            "page": i + 1,
                            "doc_type": "pptx_image",
                            "media_type": "image",
                            "media_path": img_path,
                        }
                    )
                except Exception as e:
                    logger.warning(f"Failed to extract PPTX image on slide {i+1}: {e}")

        cleaned = clean_text("\n".join(slide_text_parts))
        if cleaned:
            docs.append(
                {
                    "text": cleaned,
                    "source": filepath,
                    "page": i + 1,
                    "doc_type": "pptx",
                    "media_type": "text",
                    "media_path": None,
                }
            )

    return docs


def load_tabular_data(filepath: str) -> list[dict]:
    """Load a CSV / Excel file — each row becomes one chunk."""
    if _file_too_large(filepath):
        return []
    try:
        import pandas as pd
        df = pd.read_csv(filepath) if filepath.endswith(".csv") else pd.read_excel(filepath)
    except Exception as e:
        logger.error(f"Failed to load tabular data {filepath}: {e}")
        return []

    docs: list[dict] = []
    for index, row in df.iterrows():
        row_text = ", ".join(
            f"{col}: {val}" for col, val in row.items() if pd.notna(val)
        )
        cleaned = clean_text(row_text)
        if cleaned:
            docs.append(
                {
                    "text": cleaned,
                    "source": filepath,
                    "page": int(index) + 1,
                    "doc_type": "tabular",
                }
            )
    return docs


async def load_audio_video(filepath: str) -> list[dict]:
    """
    Transcribe an audio/video file using the OpenAI Whisper API.

    Requires API_HOST=openai or API_HOST=github and a vision-capable client.

    Returns an empty list if the client is unavailable or transcription fails.
    """
    from .config import get_settings, get_llm_client
    settings = get_settings()

    if settings.api_host not in ("openai", "github"):
        logger.warning(
            f"Whisper transcription requires API_HOST=openai or github "
            f"(current: {settings.api_host}). Skipping {filepath}."
        )
        return []

    client = get_llm_client()
    if client is None:
        logger.warning(f"LLM client unavailable — skipping audio file {filepath}.")
        return []

    if _file_too_large(filepath, max_bytes=25 * 1024 * 1024):  
        return []

    logger.info(f"Transcribing audio/video: {filepath}")
    try:
        with open(filepath, "rb") as audio_file:
            transcript = await client.audio.transcriptions.create(  
                model="whisper-1",
                file=audio_file,
                response_format="text",
            )
        cleaned = clean_text(transcript)
        if cleaned:
            return [
                {
                    "text": cleaned,
                    "source": filepath,
                    "page": None,
                    "doc_type": "media_transcript",
                }
            ]
    except Exception as e:
        logger.error(f"Failed to transcribe {filepath}: {e}")
    return []


def load_webpage(url: str) -> list[dict]:
    """Fetch and extract text from a public URL."""
    try:
        response = requests.get(url, timeout=10)
        response.raise_for_status()
        soup = BeautifulSoup(response.content, "html.parser")
        for tag in soup(["script", "style"]):
            tag.extract()
        text = soup.get_text(separator=" ", strip=True)
        cleaned = clean_text(text)
        if cleaned:
            return [{"text": cleaned, "source": url, "page": None, "doc_type": "webpage"}]
    except Exception as e:
        logger.error(f"Failed to load webpage {url}: {e}")
    return []


def load_unstructured(filepath: str) -> list[dict]:
    """Fallback loader using the `unstructured` library for unsupported formats."""
    if _file_too_large(filepath):
        return []
    try:
        from unstructured.partition.auto import partition
        elements = partition(filename=filepath)
        text = "\n\n".join(str(el) for el in elements)
        cleaned = clean_text(text)
        if cleaned:
            return [
                {
                    "text": cleaned,
                    "source": filepath,
                    "page": None,
                    "doc_type": "unstructured",
                }
            ]
    except Exception as e:
        logger.error(f"Unstructured fallback failed for {filepath}: {e}")
    return []





def load_documents_from_folder(folder_path: str | None = None) -> list[dict]:
    """
    Recursively load all supported documents from a directory.

    Dispatches to the appropriate format-specific loader based on file extension.
    Unsupported formats fall through to the ``unstructured`` fallback loader.

    Note: Audio/video files are synchronous in this function signature for
    compatibility with the synchronous incremental indexer caller. Audio files
    are skipped here and should be handled by an async caller if needed.

    Args:
        folder_path: Directory to scan. Defaults to ``<repo_root>/data/``.

    Returns:
        List of document chunk dicts ready for further processing.
    """
    if folder_path is None:
        base = os.path.dirname(os.path.dirname(os.path.dirname(os.path.abspath(__file__))))
        folder_path = os.path.join(base, "data")

    if not os.path.exists(folder_path):
        os.makedirs(folder_path, exist_ok=True)
        return []

    logger.info(f"Loading documents from: {folder_path}")
    documents: list[dict] = []

    for filename in os.listdir(folder_path):
        filepath = os.path.join(folder_path, filename)
        if not os.path.isfile(filepath):
            continue

        ext = os.path.splitext(filename)[1].lower()

        if ext == ".txt":
            documents.extend(load_txt(filepath))
        elif ext == ".pdf":
            documents.extend(load_pdf(filepath))
        elif ext == ".docx":
            documents.extend(load_docx(filepath))
        elif ext == ".pptx":
            documents.extend(load_pptx(filepath))
        elif ext in (".csv", ".xls", ".xlsx"):
            documents.extend(load_tabular_data(filepath))
        elif ext in (".html", ".htm"):
            try:
                with open(filepath, "r", encoding="utf-8", errors="replace") as f:
                    soup = BeautifulSoup(f.read(), "html.parser")
                    for tag in soup(["script", "style"]):
                        tag.extract()
                    text = soup.get_text(separator=" ", strip=True)
                    cleaned = clean_text(text)
                    if cleaned:
                        documents.append(
                            {
                                "text": cleaned,
                                "source": filepath,
                                "page": None,
                                "doc_type": "local_html",
                            }
                        )
            except Exception as e:
                logger.error(f"Failed to process HTML {filepath}: {e}")
        elif ext in (".mp3", ".mp4", ".mpeg", ".mpga", ".m4a", ".wav", ".webm"):
            logger.info(
                f"Audio/video file found: {filename}. "
                "Call load_audio_video() in an async context to transcribe."
            )
        else:
            documents.extend(load_unstructured(filepath))

    logger.info(f"Loaded {len(documents)} document chunks from {folder_path}.")
    return documents
